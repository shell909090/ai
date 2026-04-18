from pathlib import Path

from ..core.base import Extractor
from ..core.registry import registry


@registry.register("application/vnd.openxmlformats-officedocument.wordprocessingml.document")
class NativeDocxExtractor(Extractor):
    """Extract text from .docx via python-docx (paragraphs and tables)."""

    name = "python_docx"
    priority = 18

    def available(self) -> bool:
        """Check that python-docx is installed."""
        try:
            import docx  # noqa: F401
        except ImportError:
            return False
        return True

    def extract(self, path: Path) -> str:
        """Yield all paragraph and table-cell text joined by newlines."""
        import docx
        from docx.oxml.ns import qn

        doc = docx.Document(str(path))
        lines: list[str] = []
        for block in doc.element.body:
            tag = block.tag.split("}")[-1]
            if tag == "p":
                lines.append("".join(r.text for r in block.findall(f".//{qn('w:t')}") if r.text))
            elif tag == "tbl":
                for row in block.findall(f".//{qn('w:tr')}"):
                    cells = [
                        "".join(t.text for t in row_cell.findall(f".//{qn('w:t')}") if t.text)
                        for row_cell in row.findall(f".//{qn('w:tc')}")
                    ]
                    lines.append("\t".join(cells))
        return "\n".join(lines)


@registry.register("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
class NativeXlsxExtractor(Extractor):
    """Extract text from .xlsx via openpyxl (all sheets, row by row)."""

    name = "openpyxl"
    priority = 18

    def available(self) -> bool:
        """Check that openpyxl is installed."""
        try:
            import openpyxl  # noqa: F401
        except ImportError:
            return False
        return True

    def extract(self, path: Path) -> str:
        """Tab-separate cells, newline-separate rows, blank-line-separate sheets."""
        import openpyxl

        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        try:
            sheets: list[str] = []
            for ws in wb.worksheets:
                rows = []
                for row in ws.iter_rows(values_only=True):
                    rows.append("\t".join("" if v is None else str(v) for v in row))
                sheets.append("\n".join(rows))
            return "\n\n".join(sheets)
        finally:
            wb.close()


@registry.register("application/vnd.openxmlformats-officedocument.presentationml.presentation")
class NativePptxExtractor(Extractor):
    """Extract text from .pptx via python-pptx (all slides, all text frames)."""

    name = "python_pptx"
    priority = 18

    def available(self) -> bool:
        """Check that python-pptx is installed."""
        try:
            import pptx  # noqa: F401
        except ImportError:
            return False
        return True

    def extract(self, path: Path) -> str:
        """Collect text from every shape on every slide, joined by newlines."""
        from pptx import Presentation

        prs = Presentation(str(path))
        lines: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs)
                        if text:
                            lines.append(text)
        return "\n".join(lines)
